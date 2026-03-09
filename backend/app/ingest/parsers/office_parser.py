import asyncio
import logging
import unicodedata
from pathlib import Path

from .models import ParseResult

logger = logging.getLogger("litvault.parser.office")


def parse_docx(path: Path) -> ParseResult:
    """Parse a DOCX file extracting text from paragraphs and tables."""
    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(str(path))

        parts: list[str] = []

        # Extract paragraph text
        for para in doc.paragraphs:
            if para.text:
                parts.append(para.text)

        # Extract table cell text
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        parts.append(cell.text)

        text = unicodedata.normalize("NFC", "\n".join(parts))
        return ParseResult(
            text=text,
            page_count=0,  # DOCX has no page concept
            has_text=len(text.strip()) > 0,
        )

    except Exception as exc:
        logger.error("Failed to parse DOCX %s: %s", path, exc)
        return ParseResult(error=str(exc))


def parse_pptx(path: Path) -> ParseResult:
    """Parse a PPTX file extracting text from all slides and shapes."""
    try:
        from pptx import Presentation

        prs = Presentation(str(path))

        parts: list[str] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text
                    if text:
                        parts.append(text)

        text = unicodedata.normalize("NFC", "\n".join(parts))
        page_count = len(prs.slides)
        return ParseResult(
            text=text,
            page_count=page_count,
            has_text=len(text.strip()) > 0,
        )

    except Exception as exc:
        logger.error("Failed to parse PPTX %s: %s", path, exc)
        return ParseResult(error=str(exc))


async def parse_docx_async(path: Path) -> ParseResult:
    """Async wrapper for parse_docx with a 30-second timeout."""
    try:
        return await asyncio.wait_for(
            asyncio.to_thread(parse_docx, path),
            timeout=30.0,
        )
    except asyncio.TimeoutError:
        logger.error("Timeout parsing DOCX %s", path)
        return ParseResult(error="timeout")


async def parse_pptx_async(path: Path) -> ParseResult:
    """Async wrapper for parse_pptx with a 30-second timeout."""
    try:
        return await asyncio.wait_for(
            asyncio.to_thread(parse_pptx, path),
            timeout=30.0,
        )
    except asyncio.TimeoutError:
        logger.error("Timeout parsing PPTX %s", path)
        return ParseResult(error="timeout")
